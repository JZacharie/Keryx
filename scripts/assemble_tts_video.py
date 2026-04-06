"""
Assemble localized video with voice cloning (Coqui XTTS v2):
1. Remap keyframes to segments by timestamp
2. Group segments by keyframe
3. Fluify grouped text via Ollama
4. Clone voice with Coqui XTTS v2 (Mon_enregistrement_2.wav)
5. Assemble: each slide shown for its TTS duration, transitions at original segment boundaries
6. Generate PPTX
"""
import os
import json
import asyncio
import tempfile
import cv2
import requests
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
from pptx import Presentation
from pptx.util import Inches

MANIFEST       = "/app/outputs/transcripts/manifest.json"
KEYFRAMES_DIR  = "/app/outputs/keyframes/industrializing_clean"
OUTPUT_VIDEO_EN  = "/app/outputs/revoiced/industrializing_en_cloned.mp4"
OUTPUT_VIDEO_FR  = "/app/outputs/revoiced/industrializing_fr_cloned.mp4"
OUTPUT_VIDEO_ORI = "/app/outputs/revoiced/industrializing_original_clean.mp4"
ORIGINAL_AUDIO   = "/app/outputs/audio/industrializing.wav"
# CUDA / CPU toggle — set USE_CUDA=1 to enable GPU acceleration in Coqui
USE_CUDA = os.getenv("USE_CUDA", "0") == "1"
OUTPUT_PPTX    = "/app/outputs/industrializing_slides.pptx"
OLLAMA_URL     = os.getenv("OLLAMA_URL", "http://192.168.0.191:11434")
VOICE_CLONER   = os.getenv("VOICE_CLONER_URL", "http://172.18.0.4:9880")
REF_WAV        = "/app/host/Mon_enregistrement_2.wav"

def get_video_duration() -> float:
    cap = cv2.VideoCapture("/app/host/Industrializing_AI.mp4")
    fps = cap.get(cv2.CAP_PROP_FPS)
    frames = cap.get(cv2.CAP_PROP_FRAME_COUNT)
    cap.release()
    return frames / fps

def remap_keyframes(segments: list, kf_files: list, duration: float) -> list:
    n = len(kf_files)
    interval = duration / n
    for seg in segments:
        mid = (seg["start"] + seg["end"]) / 2
        seg["keyframe"] = kf_files[min(int(mid / interval), n - 1)]
    return segments

def fluify(text: str) -> str:
    prompt = (
        "Tu es un rédacteur professionnel. "
        "Réécris le texte suivant en une seule phrase fluide et naturelle en français, "
        "sans rien ajouter ni inventer, garde le sens exact. "
        "Réponds uniquement avec la phrase réécrite, sans guillemets ni explication.\n\n"
        f"{text}"
    )
    try:
        r = requests.post(f"{OLLAMA_URL}/api/generate",
                          json={"model": "llama3", "prompt": prompt, "stream": False},
                          timeout=60)
        return r.json().get("response", text).strip()
    except Exception as e:
        print(f"  [Ollama error] {e}")
        return text

def clone_voice(text: str, path: str, lang: str = "fr") -> bool:
    """Call Coqui XTTS v2 with Mon_enregistrement_2.wav as reference."""
    try:
        r = requests.get(VOICE_CLONER, params={
            "text": text,
            "language": lang,
            "speaker_wav": REF_WAV,
            "use_cuda": "1" if USE_CUDA else "0",
        }, timeout=300)
        if r.status_code == 200:
            with open(path, "wb") as f:
                f.write(r.content)
            return True
        print(f"  [Coqui error] {r.status_code}: {r.text[:100]}")
        return False
    except Exception as e:
        print(f"  [Coqui error] {e}")
        return False

def extract_elements(img_path: str, out_dir: str) -> list:
    """
    Extract main visual elements from a cleaned keyframe as transparent PNGs.
    Returns list of (png_path, x, y, w, h) in original image coordinates.
    """
    import cv2
    arr = np.array(Image.open(img_path).convert("RGB"))
    h, w = arr.shape[:2]
    gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)

    _, thresh = cv2.threshold(gray, 245, 255, cv2.THRESH_BINARY_INV)
    kernel = np.ones((15, 15), np.uint8)
    morphed = cv2.morphologyEx(thresh, cv2.MORPH_CLOSE, kernel)
    contours, _ = cv2.findContours(morphed, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    min_area = h * w * 0.05
    elements = []
    for i, cnt in enumerate(sorted(contours, key=cv2.contourArea, reverse=True)):
        if cv2.contourArea(cnt) < min_area:
            continue

        x, y, cw, ch = cv2.boundingRect(cnt)
        roi = arr[y:y+ch, x:x+cw]
        roi_gray = gray[y:y+ch, x:x+cw]

        # Build alpha: white pixels → transparent
        alpha = np.where(roi_gray > 245, 0, 255).astype(np.uint8)
        rgba = np.dstack([roi, alpha])

        png_path = os.path.join(out_dir, f"{os.path.splitext(os.path.basename(img_path))[0]}_elem{i}.png")
        Image.fromarray(rgba).save(png_path)
        os.chmod(png_path, 0o777)
        elements.append((png_path, x, y, cw, ch))

    return elements


def write_pptx(files: list):
    from pptx.util import Emu
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    elem_dir = os.path.join(os.path.dirname(OUTPUT_PPTX), "pptx_elements")
    os.makedirs(elem_dir, exist_ok=True)

    for img_path in files:
        slide = prs.slides.add_slide(blank)
        img = Image.open(img_path)
        iw, ih = img.size

        # Layer 1: full cleaned keyframe as background
        slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)

        # Layer 2+: extracted elements as transparent PNG overlays
        elements = extract_elements(img_path, elem_dir)
        for png_path, ex, ey, ew, eh in elements:
            # Scale coordinates from image pixels to slide inches
            left   = Inches(ex / iw * 13.33)
            top    = Inches(ey / ih * 7.5)
            width  = Inches(ew / iw * 13.33)
            height = Inches(eh / ih * 7.5)
            slide.shapes.add_picture(png_path, left, top, width, height)

    prs.save(OUTPUT_PPTX)
    os.chmod(OUTPUT_PPTX, 0o777)
    print(f"PPTX saved: {OUTPUT_PPTX} ({len(files)} slides with element layers)")

def assemble_original(groups: list, output_path: str):
    """Version 3: original audio track + cleaned keyframes, slide changes at segment boundaries."""
    full_audio = AudioFileClip(ORIGINAL_AUDIO)

    # Fix end times: each group ends at the start of the next (no gap)
    for i in range(len(groups) - 1):
        groups[i]['end'] = groups[i+1]['start']
    groups[-1]['end'] = full_audio.duration

    clips = []
    for i, g in enumerate(groups):
        img_path = os.path.join(KEYFRAMES_DIR, g["keyframe"])
        if not os.path.exists(img_path):
            print(f"  [SKIP] Missing: {g['keyframe']}")
            continue
        start = g["start"]
        end   = min(g["end"], full_audio.duration)
        clip  = ImageClip(img_path).set_duration(end - start).set_audio(full_audio.subclip(start, end))
        clips.append(clip)
        print(f"  [{i+1}/{len(groups)}] {g['keyframe']} [{start:.1f}→{end:.1f}s]")

    if not clips:
        print(f"No clips for {output_path}")
        return
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    final = concatenate_videoclips(clips, method="compose")
    final.write_videofile(output_path, codec="libx264", audio_codec="aac", fps=24, logger=None)
    os.chmod(output_path, 0o777)
    print(f"Video saved: {output_path}")


    """Generate cloned voice video for a given text_key ('original_en' or 'fluent')."""
    clips = []
    with tempfile.TemporaryDirectory() as tmp:
        for i, g in enumerate(groups):
            img_path = os.path.join(KEYFRAMES_DIR, g["keyframe"])
            if not os.path.exists(img_path):
                print(f"  [SKIP] Missing: {g['keyframe']}")
                continue
            text = g[text_key]
            audio_path = os.path.join(tmp, f"tts_{i}.wav")
            print(f"  [{i+1}/{len(groups)}] [{lang}] {text[:60]}...")
            if not clone_voice(text, audio_path, lang):
                print(f"  [SKIP] Clone failed for group {i}")
                continue
            audio = AudioFileClip(audio_path)
            clip = ImageClip(img_path).set_duration(audio.duration).set_audio(audio)
            clips.append(clip)

        if not clips:
            print(f"No clips for {output_path}")
            return
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        final = concatenate_videoclips(clips, method="compose")
        final.write_videofile(output_path, codec="libx264", audio_codec="aac", fps=24, logger=None)
        os.chmod(output_path, 0o777)
        print(f"Video saved: {output_path}")

def main():
    with open(MANIFEST) as f:
        segments = json.load(f)

    kf_files = sorted(os.listdir(KEYFRAMES_DIR))
    duration = get_video_duration()
    print(f"Video: {duration:.1f}s | {len(kf_files)} keyframes | {len(segments)} segments")

    segments = remap_keyframes(segments, kf_files, duration)

    # Group by keyframe, collecting both original EN and translated FR text
    groups = []
    for seg in segments:
        kf      = seg.get("keyframe")
        text_en = seg.get("text", "").strip()
        text_fr = seg.get("translated_text", "").strip()
        if not kf or not text_en:
            continue
        if groups and groups[-1]["keyframe"] == kf:
            groups[-1]["original_en"] += " " + text_en
            groups[-1]["translated"]  += " " + text_fr
            groups[-1]["end"]          = seg["end"]
        else:
            groups.append({"keyframe": kf,
                           "original_en": text_en,
                           "translated":  text_fr,
                           "start": seg["start"], "end": seg["end"]})

    print(f"{len(groups)} keyframe groups")

    # Fluify FR groups via Ollama
    print("Fluifying FR text via Ollama...")
    for i, g in enumerate(groups):
        g["fluent"] = fluify(g["translated"])
        print(f"  [{i+1}/{len(groups)}] {g['fluent'][:80]}")

    print(f"USE_CUDA={USE_CUDA}")

    # Version EN — original text, English voice cloning
    print("\n=== Generating EN version ===")
    assemble(groups, "original_en", "en", OUTPUT_VIDEO_EN)

    # Version FR — fluified translation, French voice cloning
    print("\n=== Generating FR version ===")
    assemble(groups, "fluent", "fr", OUTPUT_VIDEO_FR)

    # Version ORIGINAL — original audio + cleaned images
    print("\n=== Generating ORIGINAL audio + clean images version ===")
    assemble_original(groups, OUTPUT_VIDEO_ORI)

    # PPTX
    seen, ordered = set(), []
    for g in groups:
        p = os.path.join(KEYFRAMES_DIR, g["keyframe"])
        if p not in seen and os.path.exists(p):
            seen.add(p); ordered.append(p)
    write_pptx(ordered)

if __name__ == "__main__":
    main()
