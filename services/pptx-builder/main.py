import io
import os
import asyncio
import logging
from functools import partial
from fastapi import FastAPI
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import aioboto3
import httpx
from urllib.parse import urlparse

# Configure Logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(name)s: %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)
logger = logging.getLogger("keryx.pptx_builder")

class HealthCheckFilter(logging.Filter):
    def filter(self, record: logging.LogRecord) -> bool:
        if "/health" in record.getMessage():
            record.levelno = logging.DEBUG
            record.levelname = "DEBUG"
        return True

app = FastAPI(title="Keryx PPTX Builder")

S3_ENDPOINT = os.getenv("S3_ENDPOINT", "http://minio:9000")
S3_ACCESS_KEY = os.getenv("S3_ACCESS_KEY_ID") or os.getenv("AWS_ACCESS_KEY_ID")
S3_SECRET_KEY = os.getenv("S3_SECRET_ACCESS_KEY") or os.getenv("AWS_SECRET_ACCESS_KEY")
S3_BUCKET = os.getenv("S3_BUCKET", "keryx")

session = aioboto3.Session()

class Slide(BaseModel):
    image_url: str
    text: str = ""

class PptxRequest(BaseModel):
    job_id: str
    slides: list[Slide]
    target_path: str | None = None

async def fetch_image(url: str) -> bytes:
    parsed = urlparse(url)
    if url.startswith("file://"):
        return await asyncio.to_thread(lambda: open(url[7:], "rb").read())
    if any(h in parsed.netloc for h in ["minio", "zacharie.org", "localhost"]):
        parts = parsed.path.lstrip("/").split("/")
        bucket, key = parts[0], "/".join(parts[1:])
        async with session.client("s3", endpoint_url=S3_ENDPOINT,
                                   aws_access_key_id=S3_ACCESS_KEY,
                                   aws_secret_access_key=S3_SECRET_KEY,
                                   verify=False) as s3:
            resp = await s3.get_object(Bucket=bucket, Key=key)
            return await resp["Body"].read()
    async with httpx.AsyncClient() as client:
        resp = await client.get(url)
        return resp.content

def _build_pptx(slides_data: list[tuple[bytes, str]]) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for img_bytes, text in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        if text.strip():
            txBox = slide.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.8))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.LEFT
            p.runs[0].font.size = Pt(12)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

@app.get("/health")
async def health():
    return {"status": "ok"}

@app.post("/build")
async def build_pptx(request: PptxRequest):
    # Fetch all images concurrently
    images = await asyncio.gather(*[fetch_image(s.image_url) for s in request.slides])
    slides_data = [(img, s.text) for img, s in zip(images, request.slides)]

    # CPU-bound PPTX generation in thread pool
    pptx_bytes = await asyncio.to_thread(_build_pptx, slides_data)

    target_key = request.target_path or f"jobs/{request.job_id}/output.pptx"
    async with session.client("s3", endpoint_url=S3_ENDPOINT,
                               aws_access_key_id=S3_ACCESS_KEY,
                               aws_secret_access_key=S3_SECRET_KEY,
                               verify=False) as s3:
        await s3.put_object(
            Bucket=S3_BUCKET,
            Key=target_key,
            Body=pptx_bytes,
            ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    return {"status": "success", "url": f"{S3_ENDPOINT}/{S3_BUCKET}/{target_key}"}
if __name__ == "__main__":
    import uvicorn
    # Filter out health check access logs from uvicorn
    logging.getLogger("uvicorn.access").addFilter(HealthCheckFilter())
    uvicorn.run(app, host="0.0.0.0", port=8000, log_level="info")
